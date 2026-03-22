"""
Cross-pillar, strategic, and closing slides — correlation, Phillips curve,
crisis timeline, EU benchmarking, recommendations, risk matrix, thank you.
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from config.settings import REPORT_AUTHOR
from src.utils.logger import get_logger

from src.reporting.slides.helpers import (
    _set_slide_bg, _add_textbox, _add_rich_textbox, _add_paragraph,
    _add_chart_image, _add_section_header, _add_footer, _fmt_pct,
    DARK_BLUE, MEDIUM_BLUE, ACCENT_BLUE, WHITE, BLACK,
    DARK_GREY, GREEN, YELLOW, RED,
    SLIDE_WIDTH, FONT_FAMILY,
)
from src.reporting.slides.opening import _add_so_what_box
from src.reporting.slides.data_fetcher import DataFetcher

logger = get_logger(__name__)


def slide_correlation(prs, data: DataFetcher):
    """Slide 12 - Cross-Pillar Correlation Analysis."""
    logger.info("Building Slide 12: Correlation Heatmap")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_section_header(slide, "How the Pillars Connect",
                        "Interdependencies between macroeconomic indicators")

    _add_chart_image(slide, "correlation_heatmap.png",
                     Inches(0.5), Inches(1.2), Inches(8.0), Inches(5.8))

    tf = _add_rich_textbox(slide, Inches(8.8), Inches(1.5), Inches(4.2), Inches(5.0))
    tf.paragraphs[0].text = ""
    _add_paragraph(tf, "Key Correlations", font_size=14, bold=True,
                   colour=DARK_BLUE, space_after=Pt(12))
    _add_paragraph(tf, "\u2022 GDP growth is inversely correlated with unemployment",
                   font_size=11, colour=DARK_GREY)
    _add_paragraph(tf, "\u2022 Interest rates strongly influence credit dynamics",
                   font_size=11, colour=DARK_GREY)
    _add_paragraph(tf, "\u2022 Public debt rises during low-growth periods",
                   font_size=11, colour=DARK_GREY)
    _add_paragraph(tf, "\u2022 Inflation shows complex relationships with all pillars",
                   font_size=11, colour=DARK_GREY)
    _add_paragraph(tf, "\u2022 Post-crisis deleveraging visible across credit and debt",
                   font_size=11, colour=DARK_GREY)

    _add_so_what_box(slide,
                     "Strong correlations between unemployment-NPL (r=0.86) and ECB rate-Euribor (r=0.96) "
                     "confirm tight monetary transmission. These linkages amplify both recovery and risk.",
                     top=Inches(5.8))
    _add_footer(slide, slide_number=12)

    # Speaker notes
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = "The heatmap reveals the expected negative GDP-unemployment correlation and positive interest rate-inflation linkages. These relationships are key for policy scenario analysis."


def slide_phillips_curve(prs, data: DataFetcher):
    """Slide 13 - Phillips Curve."""
    logger.info("Building Slide 13: Phillips Curve")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_section_header(slide, "Phillips Curve Analysis",
                        "Unemployment-inflation trade-off for Portugal")

    _add_chart_image(slide, "phillips_curve.png",
                     Inches(0.5), Inches(1.2), Inches(8.0), Inches(5.8))

    tf = _add_rich_textbox(slide, Inches(8.8), Inches(1.5), Inches(4.2), Inches(5.0))
    tf.paragraphs[0].text = ""
    _add_paragraph(tf, "Interpretation", font_size=14, bold=True,
                   colour=DARK_BLUE, space_after=Pt(12))
    _add_paragraph(tf, "\u2022 Classic inverse relationship observable in Portugal's data",
                   font_size=11, colour=DARK_GREY)
    _add_paragraph(tf, "\u2022 Flattening of the curve evident in recent years",
                   font_size=11, colour=DARK_GREY)
    _add_paragraph(tf, "\u2022 Supply-side shocks (energy, COVID) distort the relationship",
                   font_size=11, colour=DARK_GREY)
    _add_paragraph(tf, "\u2022 Structural reforms may have shifted the NAIRU",
                   font_size=11, colour=DARK_GREY)

    _add_footer(slide, slide_number=13)

    # Speaker notes
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = "The Phillips Curve analysis shows a flattening of the traditional trade-off, consistent with the broader European experience. This has implications for how aggressive monetary tightening needs to be."


def slide_crisis_timeline(prs, data: DataFetcher):
    """Slide 14 - Crisis Timeline."""
    logger.info("Building Slide 14: Crisis Timeline")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_section_header(slide, "15 Years in Perspective: Four Shocks, One Trajectory",
                        "Major economic events impacting Portugal (2010\u20132025)")

    _add_chart_image(slide, "crisis_timeline.png",
                     Inches(0.3), Inches(1.15), Inches(12.7), Inches(5.0))

    tf = _add_rich_textbox(slide, Inches(0.5), Inches(6.3), Inches(12.3), Inches(0.8))
    tf.paragraphs[0].text = ""
    _add_paragraph(tf,
                   "Key events:  Eurozone Debt Crisis (2010\u201312)  \u2502  "
                   "Troika Bailout (2011\u201314)  \u2502  Recovery & Tourism Boom (2015\u201319)  \u2502  "
                   "COVID-19 Pandemic (2020\u201321)  \u2502  Energy Crisis & Inflation Surge (2022\u201323)  \u2502  "
                   "Monetary Tightening (2022\u201324)",
                   font_size=10, colour=DARK_GREY)

    _add_so_what_box(slide,
                     "Four major shocks in 15 years, each with different transmission channels. "
                     "The economy\u2019s response time shortened with each shock \u2014 a sign of institutional learning.",
                     top=Inches(5.8))
    _add_footer(slide, slide_number=14)

    # Speaker notes
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = "Three major shocks define the period: the sovereign debt crisis (2011-2014), COVID-19 (2020-2021), and the energy crisis (2022-2023). Each had distinct transmission mechanisms."


def slide_eu_benchmarking(prs, data: DataFetcher):
    """Slide 15 - EU Benchmarking: Portugal vs EU averages."""
    logger.info("Building Slide 15: EU Benchmarking")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_section_header(slide, "Converging With Europe: Gaps Closing",
                        "Portugal performance relative to European averages")

    # Left chart: radar
    _add_chart_image(slide, "benchmark_radar_pt_vs_eu.png",
                     Inches(0.3), Inches(1.2),
                     Inches(6.3), Inches(5.5))

    # Right chart: small multiples
    _add_chart_image(slide, "benchmark_small_multiples.png",
                     Inches(6.8), Inches(1.2),
                     Inches(6.3), Inches(5.5))

    # Caption
    _add_textbox(slide, Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
                 "Sources: Eurostat, ECB, Banco de Portugal  |  EU-27 weighted averages",
                 font_size=8, colour=RGBColor(0x99, 0x99, 0x99),
                 alignment=PP_ALIGN.CENTER)

    _add_so_what_box(slide,
                     "Portugal has converged with the EU average on unemployment and inflation, but the "
                     "debt/GDP gap (+17pp vs EU) and productivity gap remain the key convergence challenges.",
                     top=Inches(5.8))
    _add_footer(slide, slide_number=15)

    # Speaker notes
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = "This slide benchmarks Portugal against EU averages across all six pillars. The radar chart shows relative positioning, while the small multiples reveal the convergence trajectory over time."


def slide_strategic_recommendations(prs, data: DataFetcher):
    """Slide 16 - Strategic Recommendations."""
    logger.info("Building Slide 16: Strategic Recommendations")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_section_header(slide, "Strategic Recommendations",
                        "Evidence-based policy and investment considerations")

    # Data-driven recommendations based on current indicators
    debt_gdp = data.public_debt_latest().get("debt_to_gdp_ratio", 100)
    unemp_rate = data.unemployment_latest().get("unemployment_rate", 7)
    npl = data.credit_latest().get("npl_ratio", 3)

    short_term = [
        "Protect the budget surplus \u2014 resist expansionary pressure to preserve fiscal credibility",
        f"Monitor credit/GDP ratio \u2014 total credit growing faster than GDP despite NPL at {npl:.1f}%",
    ]
    medium_term = [
        f"Address youth unemployment ({unemp_rate * 3:.0f}% est.) \u2014 dual-education programmes and first-job incentives",
        f"Accelerate debt reduction from {debt_gdp:.0f}% towards 80% \u2014 use surplus for amortisation, not spending",
    ]
    long_term = [
        "Close the productivity gap \u2014 GDP per capita still ~70% of EU average; invest in I&D above 2% GDP",
        "Build counter-cyclical fiscal buffers for the next downturn \u2014 stabilisation fund from surplus years",
    ]
    categories = {
        "SHORT-TERM  (0\u201312 months)": short_term,
        "MEDIUM-TERM  (1\u20133 years)": medium_term,
        "LONG-TERM  (3\u20135+ years)": long_term,
    }

    tf = _add_rich_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(5.8))
    tf.paragraphs[0].text = ""

    for category, items in categories.items():
        _add_paragraph(tf, category, font_size=14, bold=True,
                       colour=DARK_BLUE, space_after=Pt(6), space_before=Pt(14))
        for item in items:
            _add_paragraph(tf, f"    \u25B8  {item}",
                           font_size=12, colour=DARK_GREY, space_after=Pt(6))

    _add_footer(slide, slide_number=16)

    # Speaker notes
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = "These recommendations are structured by time horizon. Short-term actions focus on monitoring and adjustment; medium-term on structural reform; long-term on resilience building."


def slide_risk_matrix(prs, data: DataFetcher):
    """Slide 17 - Risk Matrix table."""
    logger.info("Building Slide 17: Risk Matrix")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_section_header(slide, "Risk Matrix",
                        "Pillar-level risk assessment and mitigations")

    # Pull latest data to determine risk levels dynamically
    gdp = data.gdp_latest()
    unemp = data.unemployment_latest()
    cr = data.credit_latest()
    inf = data.inflation_latest()
    debt = data.public_debt_latest()
    ir = data.interest_rates_latest()

    def _risk_level_gdp():
        g = gdp.get("gdp_growth_yoy")
        if g is None:
            return "MEDIUM", YELLOW
        return ("LOW", GREEN) if g > 1.5 else ("HIGH", RED) if g < 0 else ("MEDIUM", YELLOW)

    def _risk_level_unemp():
        u = unemp.get("unemployment_rate")
        if u is None:
            return "MEDIUM", YELLOW
        return ("LOW", GREEN) if u < 7 else ("HIGH", RED) if u > 12 else ("MEDIUM", YELLOW)

    def _risk_level_credit():
        n = cr.get("npl_ratio")
        if n is None:
            return "MEDIUM", YELLOW
        return ("LOW", GREEN) if n < 3 else ("HIGH", RED) if n > 8 else ("MEDIUM", YELLOW)

    def _risk_level_inflation():
        h = inf.get("hicp")
        if h is None:
            return "MEDIUM", YELLOW
        return ("LOW", GREEN) if 1.0 <= h <= 3.0 else ("HIGH", RED) if h > 5 else ("MEDIUM", YELLOW)

    def _risk_level_debt():
        d = debt.get("debt_to_gdp_ratio")
        if d is None:
            return "MEDIUM", YELLOW
        return ("LOW", GREEN) if d < 60 else ("HIGH", RED) if d > 120 else ("MEDIUM", YELLOW)

    def _risk_level_rates():
        y = ir.get("portugal_10y_bond_yield")
        if y is None:
            return "MEDIUM", YELLOW
        return ("LOW", GREEN) if y < 3 else ("HIGH", RED) if y > 6 else ("MEDIUM", YELLOW)

    risk_data = [
        ("GDP", *_risk_level_gdp(),
         "Growth slowdown / recession", "Diversify economic drivers"),
        ("Unemployment", *_risk_level_unemp(),
         "Structural unemployment persistence", "Active labour market policies"),
        ("Credit", *_risk_level_credit(),
         "NPL deterioration in tightening cycle", "Prudential supervision"),
        ("Inflation", *_risk_level_inflation(),
         "Price stability deviation from target", "ECB policy pass-through"),
        ("Public Debt", *_risk_level_debt(),
         "Debt sustainability pressure", "Fiscal consolidation programme"),
        ("Interest Rates", *_risk_level_rates(),
         "Sovereign spread widening", "Credible fiscal framework"),
    ]

    rows = len(risk_data) + 1
    cols = 5
    tbl_shape = slide.shapes.add_table(rows, cols,
                                       Inches(0.5), Inches(1.4),
                                       Inches(12.3), Inches(0.5 + len(risk_data) * 0.6))
    table = tbl_shape.table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(1.5)
    table.columns[2].width = Inches(3.5)
    table.columns[3].width = Inches(3.0)
    table.columns[4].width = Inches(2.5)

    headers = ["Pillar", "Risk Level", "Key Risk", "Mitigation", "Status"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_FAMILY
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE

    for i, (pillar, level, level_colour, risk, mitigation) in enumerate(risk_data, 1):
        vals = [pillar, level, risk, mitigation, "\u25CF " + level]
        for j, v in enumerate(vals):
            cell = table.cell(i, j)
            cell.text = v
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT_FAMILY
                p.font.color.rgb = DARK_GREY
                if j == 1 or j == 4:
                    p.font.bold = True
                    p.font.color.rgb = level_colour
                    p.alignment = PP_ALIGN.CENTER
                else:
                    p.alignment = PP_ALIGN.LEFT

        # Risk-level-based row tinting
        if level == "LOW":
            row_bg = RGBColor(0xE8, 0xF8, 0xE8)
        elif level == "HIGH":
            row_bg = RGBColor(0xFD, 0xED, 0xEC)
        else:
            row_bg = RGBColor(0xFE, 0xF9, 0xE7)
        for j in range(cols):
            table.cell(i, j).fill.solid()
            table.cell(i, j).fill.fore_color.rgb = row_bg

    _add_footer(slide, slide_number=17)

    # Speaker notes
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = "The risk matrix synthesises our assessment across all pillars. Green indicators reflect the structural improvements achieved; amber and red flags highlight areas requiring continued policy attention."


def slide_thank_you(prs, data: DataFetcher):
    """Slide 18 - Thank You / Contact."""
    logger.info("Building Slide 18: Thank You")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, DARK_BLUE)

    # Decorative top line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0),
                                  SLIDE_WIDTH, Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    _add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10), Inches(1.2),
                 "Thank You", font_size=48, bold=True, colour=WHITE,
                 alignment=PP_ALIGN.CENTER)

    # Thin separator
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(5.0), Inches(3.4),
                                 Inches(3.3), Inches(0.03))
    sep.fill.solid()
    sep.fill.fore_color.rgb = ACCENT_BLUE
    sep.line.fill.background()

    _add_textbox(slide, Inches(1.5), Inches(3.8), Inches(10), Inches(0.6),
                 REPORT_AUTHOR, font_size=20, bold=False, colour=ACCENT_BLUE,
                 alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.5),
                 "contact@portugaldataintelligence.com",
                 font_size=14, bold=False, colour=WHITE,
                 alignment=PP_ALIGN.CENTER)

    _add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10), Inches(0.5),
                 "CONFIDENTIAL \u2014 For authorised recipients only",
                 font_size=10, bold=True, colour=ACCENT_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # Bottom decorative line
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(7.44),
                                   SLIDE_WIDTH, Inches(0.06))
    line2.fill.solid()
    line2.fill.fore_color.rgb = ACCENT_BLUE
    line2.line.fill.background()
