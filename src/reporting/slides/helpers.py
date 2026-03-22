"""
Slide helper utilities — shared formatting and layout functions.
"""

from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from src.reporting.shared_styles import FONT_HEADING, FONT_BODY
from src.utils.logger import get_logger
from config.settings import POWERBI_DIR

logger = get_logger(__name__)

# ---------------------------------------------------------------------------
# Colour palette (imported from shared styles)
# ---------------------------------------------------------------------------
from src.reporting.shared_styles import (
    COLOUR_PRIMARY_DARK as _PRIMARY_DARK, COLOUR_PRIMARY_MID as _PRIMARY_MID,
    COLOUR_PRIMARY_LIGHT as _PRIMARY_LIGHT, COLOUR_ACCENT_LIGHT as _ACCENT_LIGHT,
    COLOUR_BG_TINT as _BG_TINT, COLOUR_TEXT_PRIMARY as _TEXT_PRIMARY,
    COLOUR_TEXT_SECONDARY as _TEXT_SECONDARY, COLOUR_POSITIVE as _POSITIVE,
    COLOUR_NEGATIVE as _NEGATIVE, COLOUR_WARNING as _WARNING,
    COLOUR_BORDER as _BORDER, COLOUR_WHITE as _WHITE,
    pptx_rgb,
)

# Convert shared tuples to pptx RGBColor objects
COLOUR_PRIMARY_DARK = pptx_rgb(_PRIMARY_DARK)
COLOUR_PRIMARY_MID = pptx_rgb(_PRIMARY_MID)
COLOUR_PRIMARY_LIGHT = pptx_rgb(_PRIMARY_LIGHT)
COLOUR_ACCENT_LIGHT = pptx_rgb(_ACCENT_LIGHT)
COLOUR_BG_TINT = pptx_rgb(_BG_TINT)
COLOUR_TEXT_PRIMARY = pptx_rgb(_TEXT_PRIMARY)
COLOUR_TEXT_SECONDARY = pptx_rgb(_TEXT_SECONDARY)
COLOUR_POSITIVE = pptx_rgb(_POSITIVE)
COLOUR_NEGATIVE = pptx_rgb(_NEGATIVE)
COLOUR_WARNING = pptx_rgb(_WARNING)
COLOUR_BORDER = pptx_rgb(_BORDER)
COLOUR_WHITE = pptx_rgb(_WHITE)

# Semantic aliases (used throughout the presentation)
DARK_BLUE = COLOUR_PRIMARY_DARK       # Terracotta red (primary)
MEDIUM_BLUE = COLOUR_PRIMARY_MID      # Forest green (secondary)
LIGHT_BLUE = COLOUR_PRIMARY_LIGHT     # Light red (primary light)
ACCENT_BLUE = COLOUR_ACCENT_LIGHT     # Faint gold (accent light)
WHITE = COLOUR_WHITE
BLACK = COLOUR_TEXT_PRIMARY            # Near-black body text
DARK_GREY = COLOUR_TEXT_PRIMARY
LIGHT_GREY = COLOUR_BORDER            # Border/rule colour
GREEN = COLOUR_POSITIVE               # Green — positive indicators
YELLOW = COLOUR_WARNING               # Gold — warnings, highlights
RED = COLOUR_NEGATIVE                 # Red — negative indicators
VERY_LIGHT_GREY = COLOUR_BG_TINT     # Muted background

# ---------------------------------------------------------------------------
# Slide dimensions (16:9)
# ---------------------------------------------------------------------------
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ---------------------------------------------------------------------------
# Chart directory
# ---------------------------------------------------------------------------
CHARTS_DIR = POWERBI_DIR / "charts"

# ---------------------------------------------------------------------------
# Fonts
# ---------------------------------------------------------------------------
FONT_FAMILY = FONT_BODY


# ============================================================================
# Helper utilities
# ============================================================================

def _set_slide_bg(slide, colour):
    """Set a solid background colour for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = colour


def _add_textbox(slide, left, top, width, height, text, font_size=14,
                 bold=False, colour=BLACK, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_FAMILY, anchor=MSO_ANCHOR.TOP):
    """Add a simple text box and return the shape."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.paragraphs[0].alignment = alignment
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = colour
    p.font.name = font_name
    try:
        txbox.text_frame._txBody.attrib[
            "{http://schemas.openxmlformats.org/drawingml/2006/main}anchor"
        ] = "t"
    except Exception:
        pass
    return txbox


def _add_rich_textbox(slide, left, top, width, height):
    """Add a text box and return the text frame for rich paragraph building."""
    txbox = slide.shapes.add_textbox(left, top, width, height)
    tf = txbox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    return tf


def _add_paragraph(tf, text, font_size=14, bold=False, colour=BLACK,
                   alignment=PP_ALIGN.LEFT, space_after=Pt(6),
                   font_name=FONT_FAMILY, space_before=Pt(0)):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = colour
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    return p


def _add_kpi_box(slide, left, top, width, height, label, value,
                 bg_colour=DARK_BLUE, label_colour=ACCENT_BLUE,
                 value_colour=WHITE):
    """Add a KPI card (rounded rectangle with large number)."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_colour
    shape.line.fill.background()
    shape.shadow.inherit = False

    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None

    # Value (large)
    p_val = tf.paragraphs[0]
    p_val.text = str(value)
    p_val.font.size = Pt(28)
    p_val.font.bold = True
    p_val.font.color.rgb = value_colour
    p_val.font.name = FONT_HEADING
    p_val.alignment = PP_ALIGN.CENTER

    # Label (small)
    p_lbl = tf.add_paragraph()
    p_lbl.text = label
    p_lbl.font.size = Pt(9)
    p_lbl.font.bold = False
    p_lbl.font.color.rgb = label_colour
    p_lbl.font.name = FONT_BODY
    p_lbl.alignment = PP_ALIGN.CENTER
    p_lbl.space_before = Pt(2)

    return shape


def _add_chart_image(slide, chart_name, left, top, width, height):
    """Embed a chart PNG into the slide if the file exists."""
    chart_path = CHARTS_DIR / chart_name
    if chart_path.exists():
        slide.shapes.add_picture(str(chart_path), left, top, width, height)
        logger.info(f"  Embedded chart: {chart_name}")
    else:
        logger.warning(f"  Chart not found: {chart_path}")
        _add_textbox(slide, left, top, width, height,
                     f"[Chart not found: {chart_name}]",
                     font_size=12, colour=RED, alignment=PP_ALIGN.CENTER)


def _add_section_header(slide, title, subtitle=None):
    """Add a section header with title, optional subtitle and gradient bar."""
    from pptx.oxml.ns import qn
    from lxml import etree

    # Title directly on slide (no background bar)
    _add_textbox(slide, Emu(457200), Emu(181025), Emu(6831300), Emu(457200),
                 title, font_size=28, bold=True, colour=RGBColor(0x9B, 0x22, 0x26),
                 font_name=FONT_HEADING)

    if subtitle:
        _add_textbox(slide, Emu(457200), Emu(585520), Inches(10), Emu(274200),
                     subtitle, font_size=12, bold=False, colour=BLACK)

    # Gradient accent bar below header
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Emu(558900), Emu(932436),
                                 Emu(10058400), Emu(46500))
    bar.line.fill.background()

    # Apply gradient fill (green -> gold -> red)
    spPr = bar._element.spPr
    # Remove any existing fill
    for child in list(spPr):
        if child.tag.endswith('}solidFill') or child.tag.endswith('}gradFill'):
            spPr.remove(child)
    gradFill = etree.SubElement(spPr, qn('a:gradFill'))
    gsLst = etree.SubElement(gradFill, qn('a:gsLst'))
    for pos, color in [('0', '386641'), ('50000', 'D4A373'), ('100000', '9B2226')]:
        gs = etree.SubElement(gsLst, qn('a:gs'), attrib={'pos': pos})
        srgbClr = etree.SubElement(gs, qn('a:srgbClr'), attrib={'val': color})
    lin = etree.SubElement(gradFill, qn('a:lin'), attrib={'ang': '10800025', 'scaled': '0'})


def _add_footer(slide, text=None, slide_number=None):
    """Add a subtle footer at the bottom of a slide."""
    if text is None:
        if slide_number is not None:
            text = f"Portugal Data Intelligence  |  Slide {slide_number}"
        else:
            text = "Portugal Data Intelligence"
    _add_textbox(slide, Emu(457200), Emu(6446520), Emu(10972800), Emu(320040),
                 text, font_size=8, bold=False, colour=RGBColor(0x9B, 0x22, 0x26),
                 alignment=PP_ALIGN.RIGHT)


def _fmt_number(val, decimals=1, suffix="", prefix=""):
    """Format a numeric value for display."""
    if val is None:
        return "N/A"
    try:
        return f"{prefix}{val:,.{decimals}f}{suffix}"
    except (TypeError, ValueError):
        return str(val)


def _fmt_pct(val, decimals=1):
    """Format a percentage value."""
    return _fmt_number(val, decimals=decimals, suffix="%")


def _fmt_eur(val, decimals=1):
    """Format a EUR millions value."""
    if val is None:
        return "N/A"
    if abs(val) >= 1_000_000:
        return f"\u20ac{val / 1_000_000:,.{decimals}f}tn"
    if abs(val) >= 1_000:
        return f"\u20ac{val / 1_000:,.{decimals}f}bn"
    return f"\u20ac{val:,.{decimals}f}M"
