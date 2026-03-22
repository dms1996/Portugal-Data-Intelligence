"""
Opening slides — title, agenda, executive summary, scorecard.
"""

from datetime import datetime

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from config.settings import (
    REPORT_AUTHOR,
    REPORT_DATE_FORMAT,
    START_YEAR,
    END_YEAR,
)
from src.utils.logger import get_logger

from src.reporting.slides.helpers import (
    _set_slide_bg, _add_textbox, _add_rich_textbox, _add_paragraph,
    _add_kpi_box, _add_section_header, _add_footer, _fmt_pct, _fmt_eur,
    DARK_BLUE, MEDIUM_BLUE, LIGHT_BLUE, ACCENT_BLUE, WHITE, BLACK,
    DARK_GREY, GREEN, YELLOW, RED, VERY_LIGHT_GREY,
    SLIDE_WIDTH, FONT_FAMILY, FONT_HEADING,
)
from src.reporting.slides.data_fetcher import DataFetcher

logger = get_logger(__name__)


def slide_title(prs, data: DataFetcher):
    """Slide 1 - Title slide."""
    logger.info("Building Slide 1: Title")
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_slide_bg(slide, DARK_BLUE)

    # Decorative top line
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0), Inches(0),
                                  SLIDE_WIDTH, Inches(0.06))
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT_BLUE
    line.line.fill.background()

    # Title
    _add_textbox(slide, Inches(1.5), Inches(1.8), Inches(10), Inches(1.2),
                 "Portugal Data Intelligence", font_size=44, bold=True,
                 colour=WHITE, alignment=PP_ALIGN.CENTER)

    # Subtitle
    _add_textbox(slide, Inches(1.5), Inches(3.2), Inches(10), Inches(0.7),
                 f"Macroeconomic Analysis {START_YEAR}\u2013{END_YEAR}",
                 font_size=22, bold=False, colour=ACCENT_BLUE,
                 alignment=PP_ALIGN.CENTER)

    # Thin separator
    sep = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(4.5), Inches(4.1),
                                 Inches(4.3), Inches(0.03))
    sep.fill.solid()
    sep.fill.fore_color.rgb = ACCENT_BLUE
    sep.line.fill.background()

    # Date
    date_str = datetime.now().strftime(REPORT_DATE_FORMAT)
    _add_textbox(slide, Inches(1.5), Inches(4.5), Inches(10), Inches(0.5),
                 date_str, font_size=14, bold=False, colour=WHITE,
                 alignment=PP_ALIGN.CENTER)

    # Bottom decorative line
    line2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   Inches(0), Inches(7.44),
                                   SLIDE_WIDTH, Inches(0.06))
    line2.fill.solid()
    line2.fill.fore_color.rgb = ACCENT_BLUE
    line2.line.fill.background()


def slide_agenda(prs, data: DataFetcher):
    """Slide 2 - Agenda structured in 3 strategic blocks."""
    logger.info("Building Slide 2: Agenda")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_section_header(slide, "Agenda")

    blocks = {
        "I.  STABILITY ACHIEVED": [
            "Executive Summary & Scorecard",
            "Fiscal Transformation & Debt Sustainability",
            "Banking System Cleanup",
        ],
        "II.  GROWTH & EMPLOYMENT": [
            "GDP Recovery Trajectory",
            "Labour Market Turnaround",
            "Inflation: From Deflation to Target",
            "Financial Conditions",
        ],
        "III.  POSITION & OUTLOOK": [
            "EU Benchmarking & Convergence",
            "Cross-Pillar Dynamics",
            "Risk Matrix & Strategic Recommendations",
        ],
    }

    tf = _add_rich_textbox(slide, Inches(1.0), Inches(1.3), Inches(11), Inches(5.8))
    tf.paragraphs[0].text = ""

    for block_title, items in blocks.items():
        _add_paragraph(tf, block_title,
                       font_size=16, bold=True, colour=DARK_BLUE,
                       space_after=Pt(4), space_before=Pt(16))
        for item in items:
            _add_paragraph(tf, f"      \u25B8  {item}",
                           font_size=14, bold=False, colour=DARK_GREY,
                           space_after=Pt(4))

    _add_footer(slide, slide_number=2)

    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = (
        "The presentation follows a strategic narrative in 3 blocks: "
        "(I) Evidence that Portugal achieved macroeconomic stability, "
        "(II) The growth and employment story, "
        "(III) Where Portugal stands relative to the EU and what comes next. "
        "This structure mirrors the Pyramid Principle: conclusion first, then supporting evidence."
    )


def slide_executive_summary(prs, data: DataFetcher):
    """Slide 3 - Executive Summary: thesis + 3 supporting arguments."""
    logger.info("Building Slide 3: Executive Summary")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_section_header(slide, "Executive Summary",
                        f"Macroeconomic assessment {START_YEAR}\u2013{END_YEAR}")

    gdp = data.gdp_latest()
    unemp = data.unemployment_latest()
    debt = data.public_debt_latest()
    credit = data.credit_latest()
    unemp_peak = data.unemployment_peak()
    debt_peak = data.public_debt_peak()
    total_growth = data.gdp_total_growth()

    # --- Central thesis ---
    _add_textbox(slide, Inches(0.8), Inches(1.3), Inches(11.5), Inches(0.8),
                 "Portugal completed the most significant macroeconomic transformation "
                 "in the Eurozone: from bailout recipient to fiscal surplus in one decade.",
                 font_size=16, bold=True, colour=DARK_BLUE,
                 alignment=PP_ALIGN.LEFT)

    # --- 3 evidence pillars as KPI boxes ---
    box_y = Inches(2.4)
    box_h = Inches(2.8)
    box_w = Inches(3.5)
    gap = Inches(0.3)

    # Pillar 1: Fiscal Stability
    x1 = Inches(0.8)
    shape1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x1, box_y, box_w, box_h)
    shape1.fill.solid()
    shape1.fill.fore_color.rgb = RGBColor(0xF0, 0xF7, 0xF0)
    shape1.line.color.rgb = GREEN
    shape1.line.width = Pt(1.5)
    tf1 = shape1.text_frame
    tf1.word_wrap = True
    p = tf1.paragraphs[0]
    p.text = "FISCAL STABILITY"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = GREEN
    p.font.name = FONT_FAMILY
    p.alignment = PP_ALIGN.CENTER
    for label, value in [
        ("Debt/GDP", f"{_fmt_pct(debt_peak.get('debt_to_gdp_ratio'))} \u2192 {_fmt_pct(debt.get('debt_to_gdp_ratio'))}"),
        ("Budget Balance", f"{_fmt_pct(debt.get('budget_deficit'))} of GDP"),
        ("NPL Ratio", f"{_fmt_pct(credit.get('npl_ratio'))}"),
    ]:
        _add_paragraph(tf1, f"{value}", font_size=18, bold=True,
                       colour=DARK_BLUE, alignment=PP_ALIGN.CENTER, space_before=Pt(8))
        _add_paragraph(tf1, label, font_size=9, colour=DARK_GREY,
                       alignment=PP_ALIGN.CENTER)

    # Pillar 2: Growth & Employment
    x2 = x1 + box_w + gap
    shape2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x2, box_y, box_w, box_h)
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = RGBColor(0xF0, 0xF2, 0xF7)
    shape2.line.color.rgb = RGBColor(0x4A, 0x6F, 0xA5)
    shape2.line.width = Pt(1.5)
    tf2 = shape2.text_frame
    tf2.word_wrap = True
    p = tf2.paragraphs[0]
    p.text = "GROWTH & EMPLOYMENT"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x4A, 0x6F, 0xA5)
    p.font.name = FONT_FAMILY
    p.alignment = PP_ALIGN.CENTER
    for label, value in [
        ("Cumulative GDP Growth", _fmt_pct(total_growth)),
        ("Unemployment", f"{_fmt_pct(unemp_peak.get('unemployment_rate'))} \u2192 {_fmt_pct(unemp.get('unemployment_rate'))}"),
        ("GDP per Capita", _fmt_eur(gdp.get('gdp_per_capita'))),
    ]:
        _add_paragraph(tf2, f"{value}", font_size=18, bold=True,
                       colour=DARK_BLUE, alignment=PP_ALIGN.CENTER, space_before=Pt(8))
        _add_paragraph(tf2, label, font_size=9, colour=DARK_GREY,
                       alignment=PP_ALIGN.CENTER)

    # Pillar 3: Risks to Watch
    x3 = x2 + box_w + gap
    shape3 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     x3, box_y, box_w, box_h)
    shape3.fill.solid()
    shape3.fill.fore_color.rgb = RGBColor(0xFD, 0xF2, 0xF2)
    shape3.line.color.rgb = RED
    shape3.line.width = Pt(1.5)
    tf3 = shape3.text_frame
    tf3.word_wrap = True
    p = tf3.paragraphs[0]
    p.text = "RISKS TO WATCH"
    p.font.size = Pt(11)
    p.font.bold = True
    p.font.color.rgb = RED
    p.font.name = FONT_FAMILY
    p.alignment = PP_ALIGN.CENTER
    risks = [
        ("Youth Unemployment", _fmt_pct(unemp.get('youth_unemployment_rate'))),
        ("Debt/GDP vs EU Gap", "+17 pp above EU average"),
        ("Productivity Gap", "GDP/capita ~70% of EU avg"),
    ]
    for label, value in risks:
        _add_paragraph(tf3, f"{value}", font_size=18, bold=True,
                       colour=DARK_BLUE, alignment=PP_ALIGN.CENTER, space_before=Pt(8))
        _add_paragraph(tf3, label, font_size=9, colour=DARK_GREY,
                       alignment=PP_ALIGN.CENTER)

    # --- Bottom insight bar ---
    callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(0.8), Inches(5.6),
                                     Inches(11.5), Inches(0.9))
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(0xFD, 0xF6, 0xEC)
    callout.line.color.rgb = YELLOW
    callout.line.width = Pt(1)
    tf_c = callout.text_frame
    tf_c.word_wrap = True
    p = tf_c.paragraphs[0]
    run = p.add_run()
    run.text = "STRATEGIC IMPLICATION:  "
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    run.font.name = FONT_FAMILY
    run2 = p.add_run()
    run2.text = (
        "The era of easy growth (crisis recovery + inflation-driven nominal expansion) is ending. "
        "Future convergence depends on productivity gains, not cyclical tailwinds."
    )
    run2.font.size = Pt(11)
    run2.font.color.rgb = DARK_GREY
    run2.font.name = FONT_FAMILY

    _add_footer(slide, slide_number=3)

    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = (
        "TALKING POINTS: (1) Open with the thesis: Portugal's transformation is complete. "
        "(2) Walk through 3 pillars left to right: fiscal, growth, risks. "
        "(3) Emphasise that the GREEN box achievements are real but the RED box risks "
        "are the agenda for the next decade. (4) Close with the strategic implication: "
        "growth model must change from recovery-driven to productivity-driven."
    )


def _add_so_what_box(slide, text, left=Inches(0.3), top=Inches(6.0),
                     width=Inches(12.7), height=Inches(0.7)):
    """Add a 'SO WHAT?' insight callout at the bottom of a slide."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFD, 0xF6, 0xEC)
    shape.line.color.rgb = YELLOW
    shape.line.width = Pt(1)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]

    run_label = p.add_run()
    run_label.text = "SO WHAT?  "
    run_label.font.size = Pt(10)
    run_label.font.bold = True
    run_label.font.color.rgb = DARK_BLUE
    run_label.font.name = FONT_FAMILY

    run_text = p.add_run()
    run_text.text = text
    run_text.font.size = Pt(10)
    run_text.font.bold = False
    run_text.font.color.rgb = DARK_GREY
    run_text.font.name = FONT_FAMILY


def slide_scorecard(prs, data: DataFetcher):
    """Slide 4 - Scorecard: Then vs Now."""
    logger.info("Building Slide 4: Scorecard Then vs Now")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_section_header(slide, "Scorecard: Then vs Now",
                        "Key metrics at crisis peak (2013) vs latest data")

    unemp = data.unemployment_latest()
    unemp_peak = data.unemployment_peak()
    debt = data.public_debt_latest()
    debt_peak = data.public_debt_peak()
    credit = data.credit_latest()
    npl_peak = data.credit_npl_peak()
    yield_peak = data.sovereign_yield_peak()
    ir = data.interest_rates_latest()
    gdp = data.gdp_latest()

    metrics = [
        ("Unemployment",       unemp_peak.get("unemployment_rate"),      unemp.get("unemployment_rate"),       "%",  True),
        ("Debt/GDP",           debt_peak.get("debt_to_gdp_ratio"),       debt.get("debt_to_gdp_ratio"),        "%",  True),
        ("Budget Balance",     -10.3,                                     debt.get("budget_deficit"),            "%",  False),
        ("NPL Ratio",         npl_peak.get("npl_ratio"),                credit.get("npl_ratio"),               "%",  True),
        ("10Y Bond Yield",    yield_peak.get("portugal_10y_bond_yield"),ir.get("portugal_10y_bond_yield"),     "%",  True),
        ("GDP per Capita",    16000,                                     gdp.get("gdp_per_capita"),             "EUR",False),
    ]

    # Build table
    rows = len(metrics) + 1
    cols = 5
    tbl_shape = slide.shapes.add_table(rows, cols,
                                       Inches(1.0), Inches(1.5),
                                       Inches(11.0), Inches(0.5 + len(metrics) * 0.6))
    table = tbl_shape.table
    table.columns[0].width = Inches(2.5)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(2.0)
    table.columns[3].width = Inches(2.5)
    table.columns[4].width = Inches(2.0)

    headers = ["Metric", "Crisis Peak", "Latest", "Change", "Signal"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_FAMILY
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE

    for i, (name, peak, latest, unit, lower_is_better) in enumerate(metrics, start=1):
        peak_v = peak if peak is not None else 0
        latest_v = latest if latest is not None else 0

        # Metric name
        c0 = table.cell(i, 0)
        c0.text = name
        for p in c0.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.name = FONT_FAMILY
            p.font.color.rgb = DARK_GREY

        # Peak value
        c1 = table.cell(i, 1)
        c1.text = _fmt_pct(peak_v) if unit == "%" else f"\u20ac{peak_v:,.0f}"
        for p in c1.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.name = FONT_FAMILY
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = RED

        # Latest value
        c2 = table.cell(i, 2)
        c2.text = _fmt_pct(latest_v) if unit == "%" else f"\u20ac{latest_v:,.0f}"
        for p in c2.text_frame.paragraphs:
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.name = FONT_FAMILY
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = GREEN

        # Change
        c3 = table.cell(i, 3)
        diff = latest_v - peak_v
        c3.text = f"{diff:+.1f} pp" if unit == "%" else f"\u20ac{diff:+,.0f}"
        improved = (diff < 0) if lower_is_better else (diff > 0)
        for p in c3.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.name = FONT_FAMILY
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = GREEN if improved else RED

        # Signal
        c4 = table.cell(i, 4)
        c4.text = "\u2713 Improved" if improved else "\u26A0 Deteriorated"
        for p in c4.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = FONT_FAMILY
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = GREEN if improved else RED

        # Alternate row shading
        bg = VERY_LIGHT_GREY if i % 2 == 0 else WHITE
        for j in range(cols):
            table.cell(i, j).fill.solid()
            table.cell(i, j).fill.fore_color.rgb = bg

    _add_so_what_box(slide,
                     "Every major macro indicator has improved since the crisis peak. "
                     "The question is no longer 'will Portugal stabilise?' but "
                     "'can Portugal converge with core Europe?'",
                     top=Inches(5.6))

    _add_footer(slide, slide_number=4)

    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = (
        "TALKING POINTS: Walk through the table row by row. "
        "Emphasise the magnitude of change: unemployment halved, debt/GDP down 37pp, "
        "NPL from 17.5% to 2.2%. This is not incremental improvement \u2014 it is structural transformation. "
        "Transition: 'Now let us examine what drove this change.'"
    )
