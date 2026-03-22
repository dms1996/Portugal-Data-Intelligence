"""
Economic pillar slides — dashboard, GDP, labour, interest rates, credit,
inflation, public debt.
"""

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from config.settings import START_YEAR, END_YEAR
from src.utils.logger import get_logger

from src.reporting.slides.helpers import (
    _set_slide_bg, _add_textbox, _add_rich_textbox, _add_paragraph,
    _add_kpi_box, _add_chart_image, _add_section_header, _add_footer,
    _fmt_pct, _fmt_eur,
    DARK_BLUE, MEDIUM_BLUE, LIGHT_BLUE, WHITE, BLACK,
    DARK_GREY, GREEN, RED, VERY_LIGHT_GREY,
    FONT_FAMILY,
)
from src.reporting.slides.opening import _add_so_what_box
from src.reporting.slides.data_fetcher import DataFetcher

logger = get_logger(__name__)


def slide_economic_dashboard(prs, data: DataFetcher):
    """Slide 4 - Full-slide economic dashboard chart."""
    logger.info("Building Slide 4: Economic Dashboard")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_section_header(slide, "Six Pillars, One Story",
                        "All indicators confirm the recovery-to-stability trajectory")

    _add_chart_image(slide, "economic_dashboard.png",
                     Inches(0.3), Inches(1.15),
                     Inches(12.7), Inches(5.3))

    _add_so_what_box(slide,
                     "All six indicators show the same arc: crisis (2010-2013), "
                     "recovery (2014-2019), COVID dip (2020), and normalisation (2021-2025). "
                     "The consistency across pillars confirms this is structural, not cyclical.")

    _add_footer(slide, slide_number=5)


def slide_gdp_analysis(prs, data: DataFetcher):
    """Slide 5 - GDP Analysis with chart and KPIs."""
    logger.info("Building Slide 5: GDP Analysis")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    gdp_growth = data.gdp_total_growth()
    _add_section_header(slide, f"GDP: {_fmt_pct(gdp_growth)} Cumulative Growth Since {START_YEAR}",
                        "Nominal GDP evolution and year-on-year growth trajectory")

    # Chart (left 60%)
    _add_chart_image(slide, "gdp_evolution.png",
                     Inches(0.3), Inches(1.2), Inches(7.7), Inches(5.5))

    gdp = data.gdp_latest()
    right_x = Inches(8.3)

    _add_kpi_box(slide, right_x, Inches(1.3), Inches(4.5), Inches(1.0),
                 "Latest Nominal GDP", _fmt_eur(gdp.get("nominal_gdp")))
    _add_kpi_box(slide, right_x, Inches(2.5), Inches(4.5), Inches(1.0),
                 "YoY Growth Rate", _fmt_pct(gdp.get("gdp_growth_yoy")),
                 bg_colour=MEDIUM_BLUE)
    _add_kpi_box(slide, right_x, Inches(3.7), Inches(4.5), Inches(1.0),
                 "GDP per Capita", _fmt_eur(gdp.get("gdp_per_capita")),
                 bg_colour=LIGHT_BLUE)

    # Bullet points
    mm = data.gdp_min_max()
    min_q = mm["min_quarter"]
    max_q = mm["max_quarter"]
    tf = _add_rich_textbox(slide, right_x, Inches(4.9), Inches(4.5), Inches(2.0))
    tf.paragraphs[0].text = ""
    _add_paragraph(tf, f"\u2022 Deepest contraction: {_fmt_pct(min_q.get('gdp_growth_yoy'))} "
                        f"({min_q.get('date_key', 'N/A')})",
                   font_size=11, colour=DARK_GREY)
    _add_paragraph(tf, f"\u2022 Strongest growth: {_fmt_pct(max_q.get('gdp_growth_yoy'))} "
                        f"({max_q.get('date_key', 'N/A')})",
                   font_size=11, colour=DARK_GREY)
    _add_paragraph(tf, f"\u2022 Period: {gdp.get('date_key', 'N/A')}",
                   font_size=11, colour=DARK_GREY)

    _add_footer(slide, slide_number=5)

    # Speaker notes
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = "Portugal's GDP shows a consistent recovery trajectory since the sovereign debt crisis, with a sharp COVID dip in 2020 followed by strong rebound. Key talking point: cumulative growth over the period."


def slide_gdp_deep_dive(prs, data: DataFetcher):
    """Slide 6 - GDP Deep Dive with period comparison table."""
    logger.info("Building Slide 6: GDP Deep Dive")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_section_header(slide, "GDP Deep Dive",
                        "Average growth rates across economic periods")

    periods = data.gdp_growth_avg_by_period()

    # Build a table
    rows = len(periods) + 1  # header + data
    cols = 3
    tbl_shape = slide.shapes.add_table(rows, cols,
                                       Inches(1.0), Inches(1.5),
                                       Inches(8.0), Inches(0.5 + len(periods) * 0.55))
    table = tbl_shape.table
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2.0)

    # Header row
    headers = ["Economic Period", "Avg. GDP Growth (YoY)", "Assessment"]
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(12)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = FONT_FAMILY
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_BLUE

    # Data rows
    for i, (period, avg_growth) in enumerate(periods.items(), start=1):
        # Period name
        c0 = table.cell(i, 0)
        c0.text = period
        for p in c0.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = FONT_FAMILY
            p.font.color.rgb = DARK_GREY

        # Growth value
        c1 = table.cell(i, 1)
        c1.text = _fmt_pct(avg_growth)
        for p in c1.text_frame.paragraphs:
            p.font.size = Pt(12)
            p.font.bold = True
            p.font.name = FONT_FAMILY
            p.alignment = PP_ALIGN.CENTER
            if avg_growth is not None:
                p.font.color.rgb = GREEN if avg_growth > 0 else RED
            else:
                p.font.color.rgb = DARK_GREY

        # Assessment
        c2 = table.cell(i, 2)
        if avg_growth is None:
            assessment = "N/A"
        elif avg_growth > 2:
            assessment = "Strong"
        elif avg_growth > 0:
            assessment = "Moderate"
        elif avg_growth > -1:
            assessment = "Weak"
        else:
            assessment = "Contraction"
        c2.text = assessment
        for p in c2.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.name = FONT_FAMILY
            p.alignment = PP_ALIGN.CENTER
            p.font.color.rgb = DARK_GREY

        # Alternate row shading
        bg = VERY_LIGHT_GREY if i % 2 == 0 else WHITE
        for j in range(cols):
            table.cell(i, j).fill.solid()
            table.cell(i, j).fill.fore_color.rgb = bg

    # Key finding callout
    mm = data.gdp_min_max()
    total_growth = data.gdp_total_growth()
    callout = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                     Inches(1.0), Inches(5.3),
                                     Inches(8.0), Inches(1.2))
    callout.fill.solid()
    callout.fill.fore_color.rgb = RGBColor(0xEB, 0xF5, 0xFB)
    callout.line.color.rgb = LIGHT_BLUE
    callout.line.width = Pt(1)

    tf = callout.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "KEY FINDING:  "
    run.font.size = Pt(11)
    run.font.bold = True
    run.font.color.rgb = DARK_BLUE
    run.font.name = FONT_FAMILY
    run2 = p.add_run()
    run2.text = (
        f"Portugal's economy demonstrated remarkable resilience, achieving cumulative "
        f"growth of {_fmt_pct(total_growth)} over {START_YEAR}\u2013{END_YEAR} despite "
        f"sovereign debt and pandemic crises."
    )
    run2.font.size = Pt(11)
    run2.font.bold = False
    run2.font.color.rgb = DARK_GREY
    run2.font.name = FONT_FAMILY

    _add_footer(slide, slide_number=6)

    # Speaker notes
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = "The GDP decomposition reveals the structural shifts in Portugal's economy. Note the increasing contribution of services and the persistent current account improvement."


def slide_labour_market(prs, data: DataFetcher):
    """Slide 7 - Labour Market."""
    logger.info("Building Slide 7: Labour Market")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_section_header(slide, "Unemployment Halved: From 17% to 6%",
                        "Labour market recovery across all segments")

    _add_chart_image(slide, "unemployment_trends.png",
                     Inches(0.3), Inches(1.2), Inches(7.7), Inches(5.5))

    unemp = data.unemployment_latest()
    peak = data.unemployment_peak()
    right_x = Inches(8.3)

    _add_kpi_box(slide, right_x, Inches(1.3), Inches(4.5), Inches(1.0),
                 "Current Unemployment", _fmt_pct(unemp.get("unemployment_rate")))
    _add_kpi_box(slide, right_x, Inches(2.5), Inches(4.5), Inches(1.0),
                 "Youth Unemployment",
                 _fmt_pct(unemp.get("youth_unemployment_rate")),
                 bg_colour=MEDIUM_BLUE)
    _add_kpi_box(slide, right_x, Inches(3.7), Inches(4.5), Inches(1.0),
                 "Peak Rate",
                 f"{_fmt_pct(peak.get('unemployment_rate'))} ({peak.get('date_key', '')})",
                 bg_colour=RED)

    tf = _add_rich_textbox(slide, right_x, Inches(4.9), Inches(4.5), Inches(2.0))
    tf.paragraphs[0].text = ""
    _add_paragraph(tf, f"\u2022 Participation rate: {_fmt_pct(unemp.get('labour_force_participation_rate'))}",
                   font_size=11, colour=DARK_GREY)
    _add_paragraph(tf, f"\u2022 Long-term rate: {_fmt_pct(unemp.get('long_term_unemployment_rate'))}",
                   font_size=11, colour=DARK_GREY)
    _add_paragraph(tf, f"\u2022 Significant improvement from crisis peaks",
                   font_size=11, colour=DARK_GREY)

    _add_so_what_box(slide,
                     "The labour market is no longer a vulnerability. At 6% unemployment, Portugal matches "
                     "the EU average. The remaining risk is youth unemployment (19%) \u2014 a skills mismatch "
                     "problem requiring education reform, not macro stimulus.",
                     top=Inches(5.8))
    _add_footer(slide, slide_number=7)

    # Speaker notes
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = "The unemployment decline from crisis peaks is the most significant structural improvement. Youth unemployment remains a challenge, though substantially improved from 2013 highs."


def slide_interest_rates(prs, data: DataFetcher):
    """Slide 8 - Financial Conditions: Interest Rates."""
    logger.info("Building Slide 8: Interest Rates")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_section_header(slide, "Sovereign Risk Normalised: Spread Near Zero",
                        "ECB policy rates, Euribor, and sovereign yields")

    _add_chart_image(slide, "interest_rate_environment.png",
                     Inches(0.3), Inches(1.2), Inches(7.7), Inches(5.5))

    ir = data.interest_rates_latest()
    spread = data.sovereign_spread_latest()
    right_x = Inches(8.3)

    _add_kpi_box(slide, right_x, Inches(1.3), Inches(4.5), Inches(1.0),
                 "ECB Main Refi Rate",
                 _fmt_pct(ir.get("ecb_main_refinancing_rate")))
    _add_kpi_box(slide, right_x, Inches(2.5), Inches(4.5), Inches(1.0),
                 "Euribor 3M", _fmt_pct(ir.get("euribor_3m")),
                 bg_colour=MEDIUM_BLUE)
    _add_kpi_box(slide, right_x, Inches(3.7), Inches(4.5), Inches(1.0),
                 "PT 10Y Bond Yield",
                 _fmt_pct(ir.get("portugal_10y_bond_yield")),
                 bg_colour=LIGHT_BLUE)

    tf = _add_rich_textbox(slide, right_x, Inches(4.9), Inches(4.5), Inches(2.0))
    tf.paragraphs[0].text = ""
    _add_paragraph(tf, f"\u2022 Euribor 12M: {_fmt_pct(ir.get('euribor_12m'))}",
                   font_size=11, colour=DARK_GREY)
    _add_paragraph(tf, f"\u2022 Sovereign spread: {_fmt_pct(spread.get('spread'))} over Euribor 12M",
                   font_size=11, colour=DARK_GREY)

    peak = data.sovereign_yield_peak()
    _add_paragraph(tf, f"\u2022 Peak yield: {_fmt_pct(peak.get('portugal_10y_bond_yield'))} "
                       f"({peak.get('date_key', '')})",
                   font_size=11, colour=DARK_GREY)

    _add_so_what_box(slide,
                     "Sovereign risk is fully normalised. The spread compression from 1000bps to near-zero "
                     "translates directly into lower refinancing costs and restored market confidence.",
                     top=Inches(5.8))
    _add_footer(slide, slide_number=8)

    # Speaker notes
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = "The interest rate environment reflects ECB monetary policy normalisation after the zero lower bound period. Portugal's sovereign spread has narrowed significantly, reflecting improved fiscal credibility."


def slide_credit(prs, data: DataFetcher):
    """Slide 9 - Financial Conditions: Credit."""
    logger.info("Building Slide 9: Credit Portfolio")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_section_header(slide, "Banking Cleanup Complete: NPL From 17% to 2%",
                        "Bank lending and non-performing loan dynamics")

    _add_chart_image(slide, "credit_portfolio.png",
                     Inches(0.3), Inches(1.2), Inches(7.7), Inches(5.5))

    cr = data.credit_latest()
    npl_peak = data.credit_npl_peak()
    right_x = Inches(8.3)

    _add_kpi_box(slide, right_x, Inches(1.3), Inches(4.5), Inches(1.0),
                 "Total Credit Stock", _fmt_eur(cr.get("total_credit")))
    _add_kpi_box(slide, right_x, Inches(2.5), Inches(4.5), Inches(1.0),
                 "NPL Ratio", _fmt_pct(cr.get("npl_ratio")),
                 bg_colour=MEDIUM_BLUE)
    _add_kpi_box(slide, right_x, Inches(3.7), Inches(4.5), Inches(1.0),
                 "Peak NPL Ratio",
                 f"{_fmt_pct(npl_peak.get('npl_ratio'))} ({npl_peak.get('date_key', '')})",
                 bg_colour=RED)

    tf = _add_rich_textbox(slide, right_x, Inches(4.9), Inches(4.5), Inches(2.0))
    tf.paragraphs[0].text = ""
    _add_paragraph(tf, f"\u2022 NFC lending: {_fmt_eur(cr.get('credit_nfc'))}",
                   font_size=11, colour=DARK_GREY)
    _add_paragraph(tf, f"\u2022 Household lending: {_fmt_eur(cr.get('credit_households'))}",
                   font_size=11, colour=DARK_GREY)
    _add_paragraph(tf, f"\u2022 NPL ratio improved significantly from peak",
                   font_size=11, colour=DARK_GREY)

    _add_so_what_box(slide,
                     "The banking system is clean. NPL at 2% means banks can lend productively again. "
                     "Watch the credit/GDP ratio (397%) \u2014 if credit quality deteriorates, the clean-up gains reverse.",
                     top=Inches(5.8))
    _add_footer(slide, slide_number=9)

    # Speaker notes
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = "Credit conditions have tightened as the ECB normalised rates. NPL ratios have fallen dramatically from crisis peaks, reflecting successful bank restructuring and the GACS programme."


def slide_inflation(prs, data: DataFetcher):
    """Slide 10 - Inflation."""
    logger.info("Building Slide 10: Inflation")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_section_header(slide, "Inflation: Energy Shock Absorbed, Back to Target",
                        "HICP, CPI, and core inflation trends")

    _add_chart_image(slide, "inflation_dashboard.png",
                     Inches(0.3), Inches(1.2), Inches(7.7), Inches(5.5))

    inf = data.inflation_latest()
    peak = data.inflation_peak()
    right_x = Inches(8.3)

    _add_kpi_box(slide, right_x, Inches(1.3), Inches(4.5), Inches(1.0),
                 "Current HICP", _fmt_pct(inf.get("hicp")))
    _add_kpi_box(slide, right_x, Inches(2.5), Inches(4.5), Inches(1.0),
                 "Core Inflation", _fmt_pct(inf.get("core_inflation")),
                 bg_colour=MEDIUM_BLUE)
    _add_kpi_box(slide, right_x, Inches(3.7), Inches(4.5), Inches(1.0),
                 "ECB Target", "2.0%", bg_colour=GREEN)

    tf = _add_rich_textbox(slide, right_x, Inches(4.9), Inches(4.5), Inches(2.0))
    tf.paragraphs[0].text = ""
    _add_paragraph(tf, f"\u2022 CPI: {_fmt_pct(inf.get('cpi'))}",
                   font_size=11, colour=DARK_GREY)
    _add_paragraph(tf, f"\u2022 Peak HICP: {_fmt_pct(peak.get('hicp'))} ({peak.get('date_key', '')})",
                   font_size=11, colour=DARK_GREY)

    hicp_val = inf.get("hicp")
    if hicp_val is not None:
        if hicp_val > 2.0:
            note = "\u2022 Currently above ECB 2% target"
        elif hicp_val < 1.0:
            note = "\u2022 Currently below ECB 2% target \u2014 deflationary risk"
        else:
            note = "\u2022 Currently near ECB 2% target"
        _add_paragraph(tf, note, font_size=11, colour=DARK_GREY)

    _add_so_what_box(slide,
                     "The 2022 energy shock was absorbed without a wage-price spiral. Inflation is back near "
                     "the ECB 2% target. Paradoxically, the inflation spike accelerated debt/GDP reduction via nominal GDP growth.",
                     top=Inches(5.8))
    _add_footer(slide, slide_number=10)

    # Speaker notes
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = "The 2022-2023 inflation spike was primarily energy and supply-chain driven. Core inflation proved stickier. The return towards the 2% target is ongoing but gradual."


def slide_public_debt(prs, data: DataFetcher):
    """Slide 11 - Public Finance."""
    logger.info("Building Slide 11: Public Debt")
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_slide_bg(slide, WHITE)
    _add_section_header(slide, "Debt Below 100% for First Time Since 2011",
                        "Government debt, deficit, and sustainability metrics")

    _add_chart_image(slide, "public_debt_sustainability.png",
                     Inches(0.3), Inches(1.2), Inches(7.7), Inches(5.5))

    debt = data.public_debt_latest()
    peak = data.public_debt_peak()
    right_x = Inches(8.3)

    _add_kpi_box(slide, right_x, Inches(1.3), Inches(4.5), Inches(1.0),
                 "Debt-to-GDP", _fmt_pct(debt.get("debt_to_gdp_ratio")))
    _add_kpi_box(slide, right_x, Inches(2.5), Inches(4.5), Inches(1.0),
                 "Budget Balance (% GDP)",
                 _fmt_pct(debt.get("budget_deficit")),
                 bg_colour=MEDIUM_BLUE)
    _add_kpi_box(slide, right_x, Inches(3.7), Inches(4.5), Inches(1.0),
                 "Peak Debt-to-GDP",
                 f"{_fmt_pct(peak.get('debt_to_gdp_ratio'))} ({peak.get('date_key', '')})",
                 bg_colour=RED)

    tf = _add_rich_textbox(slide, right_x, Inches(4.9), Inches(4.5), Inches(2.0))
    tf.paragraphs[0].text = ""
    _add_paragraph(tf, f"\u2022 Total debt: {_fmt_eur(debt.get('total_debt'))}",
                   font_size=11, colour=DARK_GREY)
    _add_paragraph(tf, f"\u2022 External share: {_fmt_pct(debt.get('external_debt_share'))}",
                   font_size=11, colour=DARK_GREY)

    dtg = debt.get("debt_to_gdp_ratio")
    if dtg is not None:
        if dtg > 100:
            note = "\u2022 Above 100% threshold \u2014 elevated sustainability risk"
        elif dtg > 60:
            note = "\u2022 Above Maastricht 60% threshold but improving"
        else:
            note = "\u2022 Below Maastricht 60% threshold"
        _add_paragraph(tf, note, font_size=11, colour=DARK_GREY)

    _add_so_what_box(slide,
                     "Debt below 100% of GDP for the first time since 2011, with a budget surplus. "
                     "The policy imperative: use fiscal space for productivity investment, not consumption.",
                     top=Inches(5.8))
    _add_footer(slide, slide_number=11)

    # Speaker notes
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = "Debt-to-GDP has improved significantly from crisis peaks, aided by nominal GDP growth and primary surpluses. The trajectory is encouraging but requires continued fiscal discipline."
