"""
Portugal Data Intelligence - Executive Presentation Generator
==============================================================
Generates a professional 16:9 PowerPoint presentation with embedded charts
and live data pulled from the project SQLite database.

Usage:
    python generate_presentation.py
    python generate_presentation.py --output path/to/output.pptx
"""

import argparse
import sys
from pathlib import Path
from typing import Optional

from pptx import Presentation

# ---------------------------------------------------------------------------
# Project imports
# ---------------------------------------------------------------------------
sys.path.insert(0, str(Path(__file__).resolve().parent.parent.parent))

from config.settings import DATABASE_PATH, POWERPOINT_DIR
from src.utils.logger import get_logger, log_section

from src.reporting.slides import (
    DataFetcher,
    slide_title,
    slide_agenda,
    slide_executive_summary,
    slide_scorecard,
    slide_economic_dashboard,
    slide_gdp_analysis,
    slide_gdp_deep_dive,
    slide_labour_market,
    slide_interest_rates,
    slide_credit,
    slide_inflation,
    slide_public_debt,
    slide_correlation,
    slide_phillips_curve,
    slide_crisis_timeline,
    slide_eu_benchmarking,
    slide_strategic_recommendations,
    slide_risk_matrix,
    slide_thank_you,
)
from src.reporting.slides.helpers import SLIDE_WIDTH, SLIDE_HEIGHT

logger = get_logger(__name__)


# ============================================================================
# Main orchestrator
# ============================================================================

def generate_presentation(output_path: Optional[Path] = None) -> Path:
    """
    Build the full executive presentation and save to disk.

    Parameters
    ----------
    output_path : Path, optional
        Destination .pptx file. Defaults to the project's powerpoint directory.

    Returns
    -------
    Path
        The path to the generated presentation file.
    """
    log_section(logger, "POWERPOINT PRESENTATION GENERATOR")

    if output_path is None:
        POWERPOINT_DIR.mkdir(parents=True, exist_ok=True)
        output_path = POWERPOINT_DIR / "Portugal_Data_Intelligence_Presentation.pptx"
    else:
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)

    # Initialise presentation (16:9 widescreen)
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    logger.info(f"Presentation initialised: {SLIDE_WIDTH} x {SLIDE_HEIGHT} (16:9)")

    # Connect to database
    data = DataFetcher(DATABASE_PATH)

    # Build slides in order
    slide_builders = [
        # --- OPENING ---
        slide_title,                     #  1. Cover
        slide_agenda,                    #  2. Agenda (3 strategic blocks)
        slide_executive_summary,         #  3. Thesis + 3 pillars
        slide_scorecard,                 #  4. Then vs Now scorecard
        # --- BLOCK I: STABILITY ACHIEVED ---
        slide_public_debt,               #  5. Debt below 100%
        slide_credit,                    #  6. NPL cleanup
        slide_interest_rates,            #  7. Sovereign spread normalised
        # --- BLOCK II: GROWTH & EMPLOYMENT ---
        slide_gdp_analysis,              #  8. GDP trajectory
        slide_gdp_deep_dive,             #  9. GDP by period
        slide_labour_market,             # 10. Unemployment halved
        slide_inflation,                 # 11. Inflation back to target
        # --- BLOCK III: POSITION & OUTLOOK ---
        slide_eu_benchmarking,           # 12. EU convergence
        slide_correlation,               # 13. Cross-pillar dynamics
        slide_crisis_timeline,           # 14. 15 years in perspective
        # --- BLOCK IV: STRATEGIC ---
        slide_risk_matrix,               # 15. Risk assessment
        slide_strategic_recommendations, # 16. Recommendations
        # --- APPENDIX ---
        slide_economic_dashboard,        # 17. Technical dashboard (appendix)
        slide_thank_you,                 # 18. Close
    ]

    for builder in slide_builders:
        builder(prs, data)

    # Save
    prs.save(str(output_path))
    logger.info(f"Presentation saved: {output_path}")
    logger.info(f"Total slides: {len(prs.slides)}")
    log_section(logger, "PRESENTATION GENERATION COMPLETE")

    return output_path


# ============================================================================
# CLI entry point
# ============================================================================

if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        description="Generate the Portugal Data Intelligence executive PowerPoint presentation."
    )
    parser.add_argument(
        "--output", "-o",
        type=str,
        default=None,
        help="Output path for the .pptx file (default: reports/powerpoint/Portugal_Data_Intelligence_Presentation.pptx)",
    )
    args = parser.parse_args()

    try:
        result_path = generate_presentation(args.output)
        print(f"\nPresentation generated successfully: {result_path}")
    except Exception as exc:
        logger.error(f"Presentation generation failed: {exc}", exc_info=True)
        sys.exit(1)
